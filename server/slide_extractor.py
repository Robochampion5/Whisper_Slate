"""
slide_extractor.py — Local slide/PDF text extraction pipeline (§14.2).

Design principles:
  • All extraction runs locally — no cloud API required.
  • PDF: PyMuPDF (fitz) per-page text.
  • PPTX: python-pptx per-slide shapes + speaker notes.
  • OCR: optional pytesseract fallback, gated behind OCR_ENABLED env var.
    Requires a system-level Tesseract install; the app is fully functional
    without it — scanned pages simply produce empty chunks (excluded gracefully).
  • Enrichment hook: if ENRICHMENT_PROVIDER=openai (and OPENAI_API_KEY) are set,
    each chunk's raw text is summarised into 3–5 topic keywords via the chat API
    and stored as enriched_text alongside (never replacing) raw_text.
    With no provider configured (the default) enriched_text is always None.
    The embedding is always computed from raw_text for semantic-space consistency.

No FastAPI or SQLAlchemy imports here — pure extraction logic only.
"""

from __future__ import annotations

import io
import logging
import os
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Configuration (read once at module load)
# ---------------------------------------------------------------------------

# OCR fallback: set OCR_ENABLED=true in environment to enable.
# Requires: pip install pytesseract  AND  system Tesseract (brew install tesseract).
OCR_ENABLED: bool = os.environ.get("OCR_ENABLED", "").lower() in ("1", "true", "yes")

# Minimum char count below which OCR is attempted (if enabled).
# 30 chars ≈ a short title; anything shorter is likely a scanned/image-heavy slide.
OCR_CHAR_THRESHOLD: int = int(os.environ.get("OCR_CHAR_THRESHOLD", "30"))

# Maximum file size in bytes (default 20 MB, overridable via env var).
MAX_SLIDE_FILE_BYTES: int = int(os.environ.get("MAX_SLIDE_FILE_MB", "20")) * 1024 * 1024

# Enrichment: ENRICHMENT_PROVIDER=openai + OPENAI_API_KEY → LLM keyword extraction.
# Any other value (or unset) → skip enrichment entirely.
ENRICHMENT_PROVIDER: str = os.environ.get("ENRICHMENT_PROVIDER", "").lower()
OPENAI_API_KEY: str = os.environ.get("OPENAI_API_KEY", "")

if OCR_ENABLED:
    logger.info("slide_extractor: OCR enabled (threshold=%d chars)", OCR_CHAR_THRESHOLD)
else:
    logger.info("slide_extractor: OCR disabled (set OCR_ENABLED=true to enable)")

if ENRICHMENT_PROVIDER:
    logger.info("slide_extractor: enrichment provider = %r", ENRICHMENT_PROVIDER)
else:
    logger.info("slide_extractor: enrichment disabled (set ENRICHMENT_PROVIDER + API key to enable)")


# ---------------------------------------------------------------------------
# Output type
# ---------------------------------------------------------------------------

@dataclass
class SlideChunk:
    index: int              # 0-based slide / page number
    text: str               # locally-extracted text (raw)
    source_filename: str
    char_count: int = field(init=False)
    enriched_text: Optional[str] = None  # LLM-summarised topic keywords (if enrichment configured)

    def __post_init__(self):
        self.char_count = len(self.text)


# ---------------------------------------------------------------------------
# PDF extraction (PyMuPDF)
# ---------------------------------------------------------------------------

def _extract_pdf(data: bytes, filename: str) -> list[SlideChunk]:
    import fitz  # PyMuPDF

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        raise ValueError(f"Could not open PDF '{filename}': {exc}") from exc

    if doc.needs_pass:
        raise ValueError(
            f"'{filename}' is password-protected. Please upload an unlocked PDF."
        )

    chunks: list[SlideChunk] = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text").strip()

        # OCR fallback for scanned / image-heavy pages
        if len(text) < OCR_CHAR_THRESHOLD and OCR_ENABLED:
            text = _ocr_pdf_page(page) or text

        chunks.append(SlideChunk(index=page_num, text=text, source_filename=filename))

    doc.close()
    logger.info("PDF '%s': extracted %d page(s)", filename, len(chunks))
    return chunks


def _ocr_pdf_page(page) -> str:
    """Render a PyMuPDF page to a PIL image and run Tesseract OCR."""
    try:
        import pytesseract
        from PIL import Image as PILImage

        mat = page.get_pixmap(dpi=150)
        img = PILImage.frombytes("RGB", [mat.width, mat.height], mat.samples)
        return pytesseract.image_to_string(img).strip()
    except ImportError:
        logger.warning(
            "OCR_ENABLED=true but pytesseract or Pillow not installed. "
            "Run: pip install pytesseract pillow"
        )
        return ""
    except Exception as exc:
        logger.warning("OCR failed on page: %s", exc)
        return ""


# ---------------------------------------------------------------------------
# PPTX extraction (python-pptx)
# ---------------------------------------------------------------------------

def _extract_pptx(data: bytes, filename: str) -> list[SlideChunk]:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        prs = Presentation(io.BytesIO(data))
    except PackageNotFoundError as exc:
        raise ValueError(
            f"Could not open PPTX '{filename}': file may be corrupt or not a valid PPTX. ({exc})"
        ) from exc
    except Exception as exc:
        raise ValueError(f"Could not open PPTX '{filename}': {exc}") from exc

    chunks: list[SlideChunk] = []
    for slide_num, slide in enumerate(prs.slides):
        parts: list[str] = []

        # Main slide shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)

        # Speaker notes — often carry more explicit context than bullet points (§14.2)
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            for para in notes_frame.paragraphs:
                line = para.text.strip()
                if line:
                    parts.append(line)
        except Exception:
            pass  # notes slide may not exist

        text = " ".join(parts).strip()

        # OCR fallback: PPTX slides can be image-heavy too.
        # We re-use PyMuPDF to render if available; otherwise skip.
        if len(text) < OCR_CHAR_THRESHOLD and OCR_ENABLED:
            text = _ocr_pptx_slide(slide, slide_num) or text

        chunks.append(SlideChunk(index=slide_num, text=text, source_filename=filename))

    logger.info("PPTX '%s': extracted %d slide(s)", filename, len(chunks))
    return chunks


def _ocr_pptx_slide(slide, slide_num: int) -> str:
    """
    Render a PPTX slide to an image via LibreOffice or Pillow+pptx and OCR it.
    This is best-effort — PPTX rendering without LibreOffice is complex.
    Falls back gracefully to empty string on any failure.
    """
    try:
        import pytesseract
        # python-pptx doesn't expose a render API; we'd need LibreOffice or comtypes.
        # For MVP, log a warning and return empty — the chunk will be excluded.
        logger.debug(
            "Slide %d: below OCR threshold but PPTX slide rendering not implemented. "
            "Convert to PDF first for OCR support.",
            slide_num
        )
        return ""
    except ImportError:
        return ""


# ---------------------------------------------------------------------------
# Enrichment hook (off by default)
# ---------------------------------------------------------------------------

def _enrich_chunk(raw_text: str) -> Optional[str]:
    """
    Optional: pass raw_text through an LLM to extract 3–5 topic keywords.
    Only called when ENRICHMENT_PROVIDER is configured.

    Returns the enriched keyword string, or None on any failure.
    The result is stored alongside (never replacing) raw_text — the embedding
    is always computed from raw_text for semantic-space consistency.
    """
    if ENRICHMENT_PROVIDER != "openai" or not OPENAI_API_KEY:
        return None

    try:
        from openai import OpenAI  # imported lazily so missing package doesn't break startup

        client = OpenAI(api_key=OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a topic extractor. Given lecture slide text, "
                        "return exactly 3–5 concise topic keywords or short phrases "
                        "that capture the main concepts. Respond with only the keywords, "
                        "comma-separated, no explanation."
                    ),
                },
                {"role": "user", "content": raw_text[:2000]},
            ],
            max_tokens=80,
            temperature=0.2,
        )
        return response.choices[0].message.content.strip()
    except ImportError:
        logger.warning("ENRICHMENT_PROVIDER=openai but 'openai' package not installed. pip install openai")
        return None
    except Exception as exc:
        logger.warning("Enrichment failed for chunk: %s", exc)
        return None


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_slides(data: bytes, filename: str) -> list[SlideChunk]:
    """
    Extract one SlideChunk per slide/page from a PDF or PPTX file.

    Args:
        data:     Raw file bytes.
        filename: Original filename (used to determine format and for display).

    Returns:
        List of SlideChunk.  Empty slides (text == "") are included so the
        teacher's review list shows the correct slide count; callers should
        filter to `chunk.char_count > 0` before embedding.

    Raises:
        ValueError: For unsupported formats, password-protected files, or
                    corrupt files — with a user-readable message suitable for
                    surfacing directly in an HTTP 422 response detail.
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        chunks = _extract_pdf(data, filename)
    elif ext in (".pptx", ".ppt"):
        if ext == ".ppt":
            raise ValueError(
                "Legacy .ppt format is not supported. "
                "Please save as .pptx (File → Save As in PowerPoint)."
            )
        chunks = _extract_pptx(data, filename)
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. Please upload a .pdf or .pptx file."
        )

    # Optional enrichment pass
    if ENRICHMENT_PROVIDER:
        for chunk in chunks:
            if chunk.char_count > 0:
                chunk.enriched_text = _enrich_chunk(chunk.text)

    return chunks
